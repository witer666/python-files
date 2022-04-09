from pptx import Presentation
objOpen = Presentation(r'/Users/linux/Downloads/test.pptx')
# 遍历每一页的幻灯片
for slide in objOpen.slides:
    # 遍历每一个幻灯片的形状
    for shape in slide.shapes:
        # 判断该形状是否是一个文字框架
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            #拼接文字
            print(''.join(run.text for run in paragraph.runs))
